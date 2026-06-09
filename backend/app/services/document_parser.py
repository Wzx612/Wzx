"""Multi-format document parser: PDF, DOCX, XLSX, PPTX, Markdown, code files."""

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs",
    ".cpp", ".c", ".h", ".cs", ".rb", ".php", ".swift", ".kt",
    ".sh", ".bash", ".yaml", ".yml", ".json", ".toml", ".xml",
    ".html", ".css", ".scss", ".sql", ".r",
}


def detect_file_type(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix in (".docx", ".doc"):
        return "docx"
    if suffix in (".xlsx", ".xls"):
        return "xlsx"
    if suffix in (".pptx", ".ppt"):
        return "pptx"
    if suffix in (".md", ".markdown"):
        return "markdown"
    if suffix in CODE_EXTENSIONS:
        return "code"
    return "text"


def parse_document(content: bytes, filename: str) -> tuple[str, str]:
    """Parse document bytes into plain text.

    Returns (text, file_type).
    """
    file_type = detect_file_type(filename)

    try:
        if file_type == "pdf":
            return _parse_pdf(content), file_type
        if file_type == "docx":
            return _parse_docx(content), file_type
        if file_type == "xlsx":
            return _parse_xlsx(content), file_type
        if file_type == "pptx":
            return _parse_pptx(content), file_type
        # markdown / code / text — decode as UTF-8
        return content.decode("utf-8", errors="replace"), file_type
    except Exception as exc:
        logger.error("Parse error for %s: %s", filename, exc)
        raise ValueError(f"Failed to parse {filename}: {exc}") from exc


def _parse_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[第{i}页]\n{text.strip()}")
    return "\n\n".join(pages)


def _parse_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            parts.append(f"## {text}")
        else:
            parts.append(text)
    return "\n\n".join(parts)


def _parse_xlsx(content: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheets: list[str] = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"## 工作表: {name}\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append(f"## 幻灯片 {i}\n" + "\n".join(parts))
    return "\n\n".join(slides)
